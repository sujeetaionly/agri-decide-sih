import os
import subprocess
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def generate_hq_slides():
    os.makedirs('presentation/_hq_build', exist_ok=True)
    os.makedirs('presentation', exist_ok=True)

    # 1. Slide 1 Isolated HTML
    slide1_html = """<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
* { margin: 0; padding: 0; box-sizing: border-box; }
@page { size: 1920px 1080px; margin: 0; }
body {
    width: 1920px;
    height: 1080px;
    background: #ffffff url('../../web_deck/presentation_assets/user_cleaned_slide_1.png') no-repeat center / 100% 100%;
    position: relative;
    overflow: hidden;
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Helvetica, Arial, sans-serif;
    -webkit-font-smoothing: antialiased;
}

.slide-1-app-heading {
    position: absolute;
    top: 155px;
    left: 0;
    width: 100%;
    text-align: center;
    font-family: 'Times New Roman', Times, serif;
    font-size: 54px;
    font-weight: bold;
    color: #000000;
    letter-spacing: 0.5px;
}

.slide-1-left-content {
    position: absolute;
    top: 255px;
    left: 80px;
    width: 1050px;
    font-family: Arial, Helvetica, sans-serif;
    color: #000000;
}

.slide-1-bullet-entry {
    display: flex;
    align-items: baseline;
    gap: 16px;
    line-height: 1.15;
}

.slide-1-bullet-symbol {
    font-size: 42px;
    line-height: 1;
}

.slide-1-bullet-label {
    font-size: 38px;
    font-weight: bold;
}

.slide-1-bullet-value {
    font-size: 38px;
    font-weight: normal;
}

.slide-1-team-table {
    display: flex;
    flex-direction: column;
    gap: 12px;
    margin-top: 14px;
    margin-left: 55px;
    width: 900px;
}

.slide-1-team-row {
    display: grid;
    grid-template-columns: 55px 430px 240px;
    font-size: 33px;
    font-weight: normal;
    color: #000000;
}
</style>
</head>
<body>
    <div class="slide-1-app-heading">
        FASAL-DISHA (फसल-दिशा)
    </div>

    <div class="slide-1-left-content">
        <!-- Bullet 1 -->
        <div class="slide-1-bullet-entry" style="margin-bottom: 38px;">
            <span class="slide-1-bullet-symbol">•</span>
            <div>
                <span class="slide-1-bullet-label">Problem Statement ID –</span>
                <span class="slide-1-bullet-value" style="margin-left: 14px;">PS24</span>
            </div>
        </div>

        <!-- Bullet 2 -->
        <div class="slide-1-bullet-entry" style="margin-bottom: 36px;">
            <span class="slide-1-bullet-symbol">•</span>
            <div>
                <span class="slide-1-bullet-label">Problem Statement Title-</span>
                <span class="slide-1-bullet-value" style="margin-left: 8px;">AI-Based Crop</span><br>
                <span class="slide-1-bullet-value">Recommendation Engine for Farmers</span>
            </div>
        </div>

        <!-- Bullet 3 -->
        <div class="slide-1-bullet-entry" style="margin-bottom: 38px;">
            <span class="slide-1-bullet-symbol">•</span>
            <div>
                <span class="slide-1-bullet-label">Theme-</span>
                <span class="slide-1-bullet-value" style="margin-left: 8px;">Agriculture, Food Technology &amp; Rural</span><br>
                <span class="slide-1-bullet-value">Development</span>
            </div>
        </div>

        <!-- Bullet 4 -->
        <div class="slide-1-bullet-entry" style="margin-bottom: 34px;">
            <span class="slide-1-bullet-symbol">•</span>
            <div>
                <span class="slide-1-bullet-label">PS Category-</span>
                <span class="slide-1-bullet-value" style="margin-left: 8px;">Software</span>
            </div>
        </div>

        <!-- Bullet 5: Team Details -->
        <div class="slide-1-bullet-entry" style="flex-direction: column; gap: 2px;">
            <div style="display: flex; align-items: baseline; gap: 16px;">
                <span class="slide-1-bullet-symbol">•</span>
                <span class="slide-1-bullet-label">Team Details :</span>
            </div>

            <div class="slide-1-team-table">
                <div class="slide-1-team-row">
                    <span>1.</span>
                    <span>Sujeet Sandeep Bade</span>
                    <span>2025UCP1185</span>
                </div>
                <div class="slide-1-team-row">
                    <span>2.</span>
                    <span>Krishna Maheshwari</span>
                    <span>2025UCP1422</span>
                </div>
                <div class="slide-1-team-row">
                    <span>3.</span>
                    <span>Sankirat Kaur</span>
                    <span>2025UCP1706</span>
                </div>
                <div class="slide-1-team-row">
                    <span>4.</span>
                    <span>Hemlata Chopda</span>
                    <span>2025UCP1343</span>
                </div>
                <div class="slide-1-team-row">
                    <span>5.</span>
                    <span>Arnav Goel</span>
                    <span>2025UCP1138</span>
                </div>
                <div class="slide-1-team-row">
                    <span>6.</span>
                    <span>Aman Ali</span>
                    <span>2025UCP1488</span>
                </div>
            </div>
        </div>
    </div>
</body>
</html>
"""

    # 2. Slide 7 Isolated HTML
    slide7_html = """<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<style>
* { margin: 0; padding: 0; box-sizing: border-box; }
@page { size: 1920px 1080px; margin: 0; }
body {
    width: 1920px;
    height: 1080px;
    background: #ffffff url('../../web_deck/presentation_assets/user_cleaned_slide_2.png') no-repeat center / 100% 100%;
    position: relative;
    overflow: hidden;
    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Helvetica, Arial, sans-serif;
    -webkit-font-smoothing: antialiased;
}

.team-oval {
    position: absolute;
    top: 40px;
    left: 45px;
    width: 200px;
    height: 120px;
    display: flex;
    align-items: center;
    justify-content: center;
    text-align: center;
    font-family: Arial, Helvetica, sans-serif;
    font-size: 26px;
    font-weight: bold;
    color: #000000;
    z-index: 30;
}

.title-box {
    position: absolute;
    top: 25px;
    left: 270px;
    right: 370px;
    text-align: center;
    z-index: 25;
    display: flex;
    flex-direction: column;
    align-items: center;
    gap: 8px;
}
.slide-title {
    font-family: 'Times New Roman', Times, serif;
    font-size: 38px;
    font-weight: bold;
    color: #000000;
    letter-spacing: 0.5px;
    line-height: 1.15;
}
.slide-subtitle {
    font-size: 14px;
    font-weight: 700;
    color: #14532d;
    background: #f0fdf4;
    padding: 4px 18px;
    border-radius: 20px;
    border: 1.5px solid #86efac;
    letter-spacing: 0.2px;
    white-space: nowrap;
}

.flow-container {
    position: absolute;
    top: 195px;
    left: 36px;
    right: 36px;
    height: 805px;
    display: grid;
    grid-template-columns: repeat(5, 1fr);
    gap: 16px;
    z-index: 20;
}

.step-col {
    display: flex;
    flex-direction: column;
    align-items: center;
    background: #ffffff;
    border: 1.5px solid #cbd5e1;
    border-radius: 14px;
    padding: 8px 8px 8px 8px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.06);
    height: 100%;
    position: relative;
}

.step-header-box {
    width: 100%;
    background: linear-gradient(135deg, #15803d 0%, #166534 100%);
    color: #ffffff;
    border-radius: 8px;
    padding: 5px 6px;
    display: flex;
    align-items: center;
    justify-content: center;
    gap: 6px;
    margin-bottom: 6px;
    box-shadow: 0 2px 4px rgba(22, 101, 52, 0.2);
}
.step-badge {
    background: #ffffff;
    color: #15803d;
    font-size: 11px;
    font-weight: 900;
    padding: 2px 6px;
    border-radius: 5px;
    letter-spacing: 0.5px;
}
.step-title {
    font-size: 13px;
    font-weight: 800;
    color: #ffffff;
    white-space: nowrap;
    letter-spacing: 0.2px;
}

.phone-frame {
    width: 100%;
    height: 645px;
    border: none;
    background: transparent;
    box-shadow: none;
    display: flex;
    align-items: center;
    justify-content: center;
    position: relative;
    margin-bottom: 6px;
}
.phone-img {
    width: 100%;
    height: 100%;
    object-fit: contain;
    background: transparent;
    display: block;
}

.step-desc-card {
    width: 100%;
    background: #f8fafc;
    border: 1px solid #e2e8f0;
    border-radius: 8px;
    padding: 5px 7px;
    display: flex;
    flex-direction: column;
    gap: 2px;
    text-align: left;
    flex-grow: 1;
    justify-content: center;
}
.desc-top-row {
    display: flex;
    align-items: center;
    justify-content: space-between;
}
.desc-headline {
    font-size: 12.5px;
    font-weight: 800;
    color: #0f172a;
}
.desc-metric {
    font-size: 10.5px;
    font-weight: 800;
    color: #15803d;
    background: #dcfce7;
    padding: 1px 5px;
    border-radius: 5px;
}
.desc-body {
    font-size: 11px;
    line-height: 1.3;
    color: #475569;
    font-weight: 600;
}

.footer-banner-overlay {
    position: absolute;
    bottom: 0px;
    left: 0px;
    right: 0px;
    height: 58px;
    background: #0070c0;
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 0 50px;
    color: #ffffff;
    font-family: Arial, Helvetica, sans-serif;
    font-size: 16.5px;
    font-weight: bold;
    z-index: 40;
}
</style>
</head>
<body>
    <div class="team-oval">Almost AC</div>

    <div class="title-box">
        <div class="slide-title">Screenshots of Real Working Prototype</div>
        <div class="slide-subtitle">
            12+ Indian Regional Languages ➔ 250m GIS Soil Profiling ➔ Head-to-Head Profit Engine ➔ What-If Risk Sandbox ➔ 120-Day Action Plan
        </div>
    </div>

    <div class="flow-container">
        <!-- Step 1: Home & Voice AI -->
        <div class="step-col">
            <div class="step-header-box">
                <span class="step-badge">STEP 1</span>
                <span class="step-title">Home &amp; Voice AI</span>
            </div>
            <div class="phone-frame">
                <img src="../../web_deck/presentation_assets/flow_step_1_home.png" class="phone-img" alt="Home Dashboard">
            </div>
            <div class="step-desc-card">
                <div class="desc-top-row">
                    <span class="desc-headline">All Regional Languages</span>
                    <span class="desc-metric">BHASHINI AI</span>
                </div>
                <div class="desc-body">
                    Voice guidance in 12+ Indian languages, live Mandi rate ticker &amp; 1-tap 2-min crop assessment.
                </div>
            </div>
        </div>

        <!-- Step 2: 5-Card Farm Profiler -->
        <div class="step-col">
            <div class="step-header-box">
                <span class="step-badge">STEP 2</span>
                <span class="step-title">Zero-Lab GIS Profile</span>
            </div>
            <div class="phone-frame">
                <img src="../../web_deck/presentation_assets/flow_step_2_soil.png" class="phone-img" alt="Soil Profiler">
            </div>
            <div class="step-desc-card">
                <div class="desc-top-row">
                    <span class="desc-headline">Zero-Lab Satellite GIS</span>
                    <span class="desc-metric">250m Grid</span>
                </div>
                <div class="desc-body">
                    Auto-fetches SoilGrids 250m texture &amp; pH. Factorizes owned tools (tractor −₹3.5k/ac) &amp; crop rotation.
                </div>
            </div>
        </div>

        <!-- Step 3: Head-to-Head Comparison -->
        <div class="step-col">
            <div class="step-header-box">
                <span class="step-badge">STEP 3</span>
                <span class="step-title">Head-to-Head Decision</span>
            </div>
            <div class="phone-frame">
                <img src="../../web_deck/presentation_assets/flow_step_3_comparison.png" class="phone-img" alt="Head to Head Comparison">
            </div>
            <div class="step-desc-card">
                <div class="desc-top-row">
                    <span class="desc-headline">Intended vs AI Best</span>
                    <span class="desc-metric">+₹15,016/Ac</span>
                </div>
                <div class="desc-body">
                    Farmer Cotton (₹3,887) vs AI Moong (₹18,903). Itemized CACP A₂+FL cost breakdown &amp; profit delta.
                </div>
            </div>
        </div>

        <!-- Step 4: What-If Risk Sandbox -->
        <div class="step-col">
            <div class="step-header-box">
                <span class="step-badge">STEP 4</span>
                <span class="step-title">What-If Risk Sandbox</span>
            </div>
            <div class="phone-frame">
                <img src="../../web_deck/presentation_assets/flow_step_4_whatif.png" class="phone-img" alt="What-If Risk Sandbox">
            </div>
            <div class="step-desc-card">
                <div class="desc-top-row">
                    <span class="desc-headline">Pre-Sowing Stress-Test</span>
                    <span class="desc-metric">Real-Time</span>
                </div>
                <div class="desc-body">
                    Dynamic sliders for −30% rain deficit &amp; −25% price crash. Instant sensitivity recalibration (₹19,060/ac).
                </div>
            </div>
        </div>

        <!-- Step 5: 120-Day Action Calendar -->
        <div class="step-col">
            <div class="step-header-box">
                <span class="step-badge">STEP 5</span>
                <span class="step-title">120-Day Action Plan</span>
            </div>
            <div class="phone-frame">
                <img src="../../web_deck/presentation_assets/flow_step_5_calendar.png" class="phone-img" alt="Action Calendar">
            </div>
            <div class="step-desc-card">
                <div class="desc-top-row">
                    <span class="desc-headline">Operational Guidance</span>
                    <span class="desc-metric">1-Click PDF</span>
                </div>
                <div class="desc-body">
                    Day-by-day seed treatment, weeding &amp; harvest plan + 1-Click offline printable QR slip for CSCs &amp; WhatsApp.
                </div>
            </div>
        </div>
    </div>

    <div class="footer-banner-overlay">
        <span>@SIH Idea submission- Template</span>
        <span>7</span>
    </div>
</body>
</html>
"""

    s1_path = os.path.abspath('presentation/_hq_build/slide1.html')
    s7_path = os.path.abspath('presentation/_hq_build/slide7.html')

    with open(s1_path, 'w', encoding='utf-8') as f:
        f.write(slide1_html)
    with open(s7_path, 'w', encoding='utf-8') as f:
        f.write(slide7_html)

    browser_path = r'C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe'
    img1_path = os.path.abspath('presentation/slide_1_hq.png')
    img7_path = os.path.abspath('presentation/slide_7_hq.png')

    # Render High-Resolution full-bleed PNGs (1920x1080)
    for html_file, out_img in [(s1_path, img1_path), (s7_path, img7_path)]:
        url = 'file:///' + html_file.replace('\\', '/')
        cmd = [browser_path, '--headless', '--disable-gpu', f'--screenshot={out_img}', '--window-size=1920,1080', url]
        subprocess.run(cmd)
        print(f"Rendered {out_img}: {os.path.exists(out_img)}")

    # 3. Create final.pptx
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    assets_dir = os.path.abspath("web_deck/presentation_assets")

    # ====================================================
    # SLIDE 1: High-Quality Full Bleed Image (Zero-Shift, Pixel-Perfect)
    # ====================================================
    slide1 = prs.slides.add_slide(blank_layout)
    slide1.shapes.add_picture(img1_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # ====================================================
    # SLIDE 2: PROPOSED SOLUTION
    # ====================================================
    slide2 = prs.slides.add_slide(blank_layout)
    bg2_path = os.path.join(assets_dir, "user_cleaned_slide_2.png")
    if os.path.exists(bg2_path):
        slide2.shapes.add_picture(bg2_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    s2_title_box = slide2.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf2 = s2_title_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Fasal-Disha (फसल-दिशा)"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Times New Roman"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 0, 0)

    p2_sub = tf2.add_paragraph()
    p2_sub.text = "हर खेत को मिले सही दिशा — AI-Powered Net Profit (₹/Acre) Crop Decision Platform"
    p2_sub.alignment = PP_ALIGN.CENTER
    p2_sub.font.name = "Arial"
    p2_sub.font.size = Pt(11)
    p2_sub.font.bold = True
    p2_sub.font.color.rgb = RGBColor(20, 83, 45)

    s2_left = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s2_left = s2_left.text_frame
    tf_s2_left.word_wrap = True
    p = tf_s2_left.paragraphs[0]
    p.text = "❖ IDEA / PROPOSED SOLUTION :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    p = tf_s2_left.add_paragraph()
    p.text = "86.2% of Indian smallholders operate without soil health cards and choose crops by guesswork. Fasal-Disha replaces this with GPS-based soil analysis, government CACP cost data, and real mandi prices — delivering true Net Profit (₹/Acre) forecasts in under 90 seconds."
    p.font.name = "Arial"
    p.font.size = Pt(10.5)
    p.font.color.rgb = RGBColor(30, 41, 59)
    p.space_after = Pt(6)

    s2_bullets = [
        ("Zero-Lab Soil Intelligence: ", "GPS auto-fetches SoilGrids 250m satellite data (pH, clay%, organic carbon) — eliminates the 3-4 week lab test barrier."),
        ("True Net Profit, Not Gross Yield: ", "Uses official CACP A₂+FL cost norms to show real pocket profit per acre — credits tractor (−₹3,500), sprayer (−₹800), pump (−₹600) if owned."),
        ("Crop Rotation Intelligence: ", "Penalizes monoculture (−15% yield penalty) and rewards legume rotation (+12% nitrogen fixation bonus) to prevent soil exhaustion."),
        ("Sowing Window Protection: ", "ICAR-calibrated delay decay (−0.5%/day) warns farmers about invisible yield loss from late planting and auto-suggests contingency crops.")
    ]

    for title, desc in s2_bullets:
        p = tf_s2_left.add_paragraph()
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    p = tf_s2_left.add_paragraph()
    p.space_before = Pt(8)
    p.text = "❖ WHAT MAKES IT UNIQUE :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    uvp_bullets = [
        ("Head-to-Head Crop Comparison: ", "Direct ₹/Acre profit comparison between farmer's intended crop and AI optimal recommendation."),
        ("What-If Risk Simulation: ", "Stress-test rainfall deficit (−35%) and mandi price crash (−25%) before spending capital on seeds."),
        ("Multilingual Voice AI: ", "Full voice onboarding and audio playback in 12+ regional languages for non-literate farmers."),
        ("1-Click Offline Slip: ", "Instant printable A4 advisory slip with complete crop calendar for CSCs, Kisan Kendras & WhatsApp.")
    ]

    for title, desc in uvp_bullets:
        p = tf_s2_left.add_paragraph()
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    s2_right = slide2.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s2_right = s2_right.text_frame
    tf_s2_right.word_wrap = True
    p = tf_s2_right.paragraphs[0]
    p.text = "❖ FASAL-DISHA INTERACTIVE FLOW :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    flow_steps = [
        ("1. 5-Card Farm Profile", "GPS SoilGrids 250m GIS + Texture Cards + Water Source + Owned Machinery + Previous Crop."),
        ("2. Real-Time Mandi & Cost Sync", "14,780+ APMC Mandi price trend points + CACP A₂+FL official operational input costs."),
        ("3. XGBoost Predictive Engine", "District-level yield regression (R²=0.9907) adjusted with ICAR sowing delay curve (−0.5%/day)."),
        ("4. Head-to-Head Profit Matrix", "Farmer's considered choice vs AI recommendation with exact ₹/Acre net gain."),
        ("5. Live What-If Risk Sandbox", "Rainfall sensitivity and mandi crash sliders dynamically recalibrate profit."),
        ("6. 120-Day Crop Calendar & Slip", "Phase-wise farming schedule + 1-Click printable vector advisory slip.")
    ]

    for step, desc in flow_steps:
        p = tf_s2_right.add_paragraph()
        p.space_before = Pt(5)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = step + "\n"
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(71, 85, 105)

    # ====================================================
    # SLIDE 3: TECHNICAL APPROACH
    # ====================================================
    slide3 = prs.slides.add_slide(blank_layout)
    bg3_path = os.path.join(assets_dir, "user_cleaned_slide_3.png")
    if os.path.exists(bg3_path):
        slide3.shapes.add_picture(bg3_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    s3_title_box = slide3.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf3 = s3_title_box.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "Fasal-Disha Technical Approach"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = "Times New Roman"
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(0, 0, 0)

    p3_sub = tf3.add_paragraph()
    p3_sub.text = "4-Tier Architecture • Zero-Lab GIS Ingestion • XGBoost (R²=0.9907) • 100% CACP Fidelity"
    p3_sub.alignment = PP_ALIGN.CENTER
    p3_sub.font.name = "Arial"
    p3_sub.font.size = Pt(11)
    p3_sub.font.bold = True
    p3_sub.font.color.rgb = RGBColor(20, 83, 45)

    s3_left = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s3_left = s3_left.text_frame
    tf_s3_left.word_wrap = True
    p = tf_s3_left.paragraphs[0]
    p.text = "❖ 4-TIER CLOUD-NATIVE ARCHITECTURE :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    tiers = [
        ("Tier 1: Client & Voice UI", "React 18 PWA, TypeScript, Tailwind CSS, Vite, Capacitor 8.5 Native Android Container (< 1.1 MB), Web Speech Recognition, MeitY Indic Speech TTS."),
        ("Tier 2: High-Performance Engine", "FastAPI (Python 3.11 ASGI), Pydantic v2 strict schemas, CACP A₂+FL cost computation engine, ICAR sowing window delay loss calculators."),
        ("Tier 3: Machine Learning Inference", "XGBoostRegressor (R²=0.9907, 95.9% census match) for taluka yields, Agmarknet 5-yr modal pricing forecaster (R²=0.8456)."),
        ("Tier 4: Ground-Truth Repositories", "ISRIC SoilGrids 250m GIS raster layers, MoA Agmarknet APMC mandi archives, CACP Kharif/Rabi price policy norms, PostgreSQL 16 / SQLite fallback.")
    ]

    for title, desc in tiers:
        p = tf_s3_left.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title + ": "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    s3_right = slide3.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s3_right = s3_right.text_frame
    tf_s3_right.word_wrap = True
    p = tf_s3_right.paragraphs[0]
    p.text = "❖ GROUND-TRUTH ACCURACY BENCHMARKS :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    benchmarks = [
        ("Yield Prediction: ", "XGBoost vs MoA DES Census actuals — 95.9% district yield match (R²=0.9907)."),
        ("Price Forecasting: ", "Model vs Agmarknet APMC actuals — 98.6% wholesale rate accuracy across 14,780+ mandis."),
        ("Cost Calculation: ", "100% CACP A₂+FL fidelity (with verified tractor saving −₹3,500/ac, sprayer −₹800/ac)."),
        ("Sowing Delay Loss: ", "−0.5%/day calibrated from ICAR-CRIDA field trials — 100% ICAR calibration."),
        ("Real Field Proof: ", "Wardha (Maharashtra) farmer, Cotton (monoculture) → Soybean (AI recommendation) = +₹10,370/acre (+72.3% net profit gain).")
    ]

    for title, desc in benchmarks:
        p = tf_s3_right.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.8)
        r2.font.color.rgb = RGBColor(71, 85, 105)

    # ====================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ====================================================
    slide4 = prs.slides.add_slide(blank_layout)
    bg4_path = os.path.join(assets_dir, "user_cleaned_slide_4.png")
    if os.path.exists(bg4_path):
        slide4.shapes.add_picture(bg4_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    s4_title_box = slide4.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf4 = s4_title_box.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "Feasibility, Viability & Mitigations"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.name = "Times New Roman"
    p4.font.size = Pt(28)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(0, 0, 0)

    s4_left = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s4_left = s4_left.text_frame
    tf_s4_left.word_wrap = True
    p = tf_s4_left.paragraphs[0]
    p.text = "⚖️ FEASIBILITY & VIABILITY ANALYSIS :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    feas_points = [
        ("1. Technical Readiness: ", "Production stack deployed and tested — FastAPI ASGI + React 18 PWA + Capacitor Android (< 1.1 MB bundle). Tested on entry-level 4G devices."),
        ("2. Zero-Lab Removes #1 Barrier: ", "86.2% of Indian smallholders lack lab test kits. Automated 250m GIS SoilGrids satellite ingestion eliminates the 3-4 week wait."),
        ("3. Ultra-Low Overhead: ", "< ₹0.02 operational cost per recommendation query. 100% open-source FOSS licensing."),
        ("4. Deployment Ready: ", "Interoperable open REST APIs ready for integration with PM-Kisan, KVK extension services, and AgriStack national DPI.")
    ]

    for title, desc in feas_points:
        p = tf_s4_left.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    p = tf_s4_left.add_paragraph()
    p.space_before = Pt(8)
    p.text = "🛡️ REAL-WORLD CHALLENGES & MITIGATIONS :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)

    challenges = [
        ("Extreme Weather Uncertainty: ", "What-If Risk Sandbox stress-tests rainfall deficit (−35%) and mandi price crashes (−25%) before seed sowing."),
        ("Rural 2G/3G Connectivity: ", "Offline-first PWA caching + 1-Click printable vector A4 advisory slip for CSCs and Kisan Kendras."),
        ("Low Literacy & Language: ", "Full voice navigation and audio narration in 12+ regional languages — zero reading barrier.")
    ]

    for title, desc in challenges:
        p = tf_s4_left.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(185, 28, 28)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    s4_right = slide4.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s4_right = s4_right.text_frame
    tf_s4_right.word_wrap = True
    p = tf_s4_right.paragraphs[0]
    p.text = "📊 GOVERNMENT DATA GROUND-TRUTH ACCURACY :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    table_data = [
        ("District Yield (qtl/ac)", "MoA DES Census", "95.9% Match"),
        ("Cultivation Cost (₹/ac)", "CACP A₂+FL Norms", "100% Fidelity"),
        ("Mandi Wholesale Rate", "Agmarknet APMC 5-Yr", "98.6% Accuracy"),
        ("Sowing Delay Loss", "ICAR-CRIDA Trials", "100% ICAR Calibrated"),
        ("Soil Spatial Texture", "ISRIC SoilGrids 250m", "100% Geofenced")
    ]

    for param, src, match in table_data:
        p = tf_s4_right.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"• {param:<24} | {src:<18} → "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = match
        r2.font.bold = True
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(22, 101, 52)

    # ====================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # ====================================================
    slide5 = prs.slides.add_slide(blank_layout)
    bg5_path = os.path.join(assets_dir, "user_cleaned_slide_5.png")
    if os.path.exists(bg5_path):
        slide5.shapes.add_picture(bg5_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    s5_title_box = slide5.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf5 = s5_title_box.text_frame
    tf5.word_wrap = True
    p5 = tf5.paragraphs[0]
    p5.text = "Impact & Scalability Potential"
    p5.alignment = PP_ALIGN.CENTER
    p5.font.name = "Times New Roman"
    p5.font.size = Pt(28)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(0, 0, 0)

    s5_left = slide5.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s5_left = s5_left.text_frame
    tf_s5_left.word_wrap = True
    p = tf_s5_left.paragraphs[0]
    p.text = "🌱 FARMER-LEVEL ECONOMIC IMPACT :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    impact_points = [
        ("₹10,370/Acre Profit Boost: ", "+72.3% net income jump proven in ground-truth field data (Wardha cotton monoculture → soybean rotation)."),
        ("100% Zero-Lab Accessibility: ", "Zero lab waiting time — satellite 250m GIS delivers immediate soil-tailored decisions for all 14.6 crore operational landholdings."),
        ("₹3,500/Acre Ownership Credits: ", "Recognizes owned tractors, sprayers, pumps — directly reducing computed working capital requirements."),
        ("35% Drought Risk Shield: ", "What-If risk sandbox prevents catastrophic crop failure by stress-testing weather scenarios before sowing.")
    ]

    for title, desc in impact_points:
        p = tf_s5_left.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    s5_right = slide5.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s5_right = s5_right.text_frame
    tf_s5_right.word_wrap = True
    p = tf_s5_right.paragraphs[0]
    p.text = "📈 NATIONAL MACRO & ECOLOGICAL IMPACT :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    macro_points = [
        ("National Pulse & Oilseed Security: ", "Drives adoption of pulse & oilseed rotation, reducing India's massive ₹1.38 lakh crore edible oil import bill."),
        ("Soil Health Preservation: ", "Legume nitrogen fixation (+12%) naturally restores soil microbiome and cuts chemical urea dependence by 20-30%."),
        ("Extreme Commercial Scalability: ", "FastAPI serverless microservice architecture delivers 5,000+ RPS at negligible compute cost (< ₹0.02/query).")
    ]

    for title, desc in macro_points:
        p = tf_s5_right.add_paragraph()
        p.space_before = Pt(5)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(71, 85, 105)

    # ====================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # ====================================================
    slide6 = prs.slides.add_slide(blank_layout)
    bg6_path = os.path.join(assets_dir, "user_cleaned_slide_6.png")
    if os.path.exists(bg6_path):
        slide6.shapes.add_picture(bg6_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    s6_title_box = slide6.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf6 = s6_title_box.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "Research, References & Prototype Deliverables"
    p6.alignment = PP_ALIGN.CENTER
    p6.font.name = "Times New Roman"
    p6.font.size = Pt(28)
    p6.font.bold = True
    p6.font.color.rgb = RGBColor(0, 0, 0)

    s6_left = slide6.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s6_left = s6_left.text_frame
    tf_s6_left.word_wrap = True
    p = tf_s6_left.paragraphs[0]
    p.text = "📚 GOVERNMENT & SCIENTIFIC REPOSITORIES :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    refs = [
        ("1. CACP Comprehensive Scheme (2024-25): ", "Cost of Cultivation of Principal Crops in India — official A₂+FL operation cost formulas."),
        ("2. MoA Directorate of Economics & Statistics: ", "Agricultural Statistics at a Glance — 10-year district crop yields & mandi modal price archives."),
        ("3. ISRIC World Soil Information: ", "SoilGrids 250m Global Gridded Soil Information System (pH, clay fraction, CEC, bulk density)."),
        ("4. ICAR-CRIDA Technical Bulletins: ", "Contingency Crop Planning & Sowing Window Yield Loss Models across 28 Agro-Climatic Zones.")
    ]

    for title, desc in refs:
        p = tf_s6_left.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    s6_right = slide6.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s6_right = s6_right.text_frame
    tf_s6_right.word_wrap = True
    p = tf_s6_right.paragraphs[0]
    p.text = "🚀 VERIFIED WORKING PROTOTYPE & DELIVERABLES :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    deliverables = [
        ("💻 GitHub Repository (Public): ", "github.com/sujeetaionly/agri-decide-sih (100% Tested Production Stack, 12/12 passing pytest suite)."),
        ("🎬 2-Minute Demo Video: ", "Full Android Mobile App Walkthrough demonstrating Zero-Lab Onboarding, What-If Risk Sandbox, and 1-Click PDF Advisory Slip."),
        ("📱 Native Android App: ", "Capacitor 8.5 Native APK container (< 1.1 MB) + Offline-first PWA."),
        ("🗣️ Local Language Voice: ", "MeitY Digital India BHASHINI Speech TTS engine + browser Web Speech recognition.")
    ]

    for title, desc in deliverables:
        p = tf_s6_right.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "• " + title + "\n"
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.8)
        r2.font.color.rgb = RGBColor(71, 85, 105)

    # ====================================================
    # SLIDE 7: High-Quality Full Bleed Image (Zero-Shift, Pixel-Perfect)
    # ====================================================
    slide7 = prs.slides.add_slide(blank_layout)
    slide7.shapes.add_picture(img7_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    final_path = os.path.abspath("presentation/final.pptx")
    prs.save(final_path)
    print(f"Final presentation created successfully at: {final_path}")

if __name__ == "__main__":
    generate_hq_slides()
