import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck_pptx():
    prs = Presentation()
    # Set 16:9 Widescreen dimensions (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    assets_dir = os.path.abspath("web_deck/presentation_assets")
    output_path = os.path.abspath("presentation/Fasal_Disha_SIH_2025.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    # ----------------------------------------------------
    # SLIDE 1: TITLE PAGE (Exact Digital Replica)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    bg1_path = os.path.join(assets_dir, "user_cleaned_slide_1.png")
    if os.path.exists(bg1_path):
        slide1.shapes.add_picture(bg1_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # App Heading: FASAL-DISHA (फसल-दिशा)
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.95), Inches(12.333), Inches(0.65))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0)
    tf1.margin_right = Inches(0)
    tf1.margin_top = Inches(0)
    tf1.margin_bottom = Inches(0)
    p1 = tf1.paragraphs[0]
    p1.text = "FASAL-DISHA (फसल-दिशा)"
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "Times New Roman"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(0, 0, 0)

    # Left Content Container
    left_box = slide1.shapes.add_textbox(Inches(0.55), Inches(1.72), Inches(6.8), Inches(5.3))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = Inches(0)
    tf_left.margin_right = Inches(0)
    tf_left.margin_top = Inches(0)
    tf_left.margin_bottom = Inches(0)

    items = [
        ("Problem Statement ID –", "  PS24", 8),
        ("Problem Statement Title-", " AI-Based Crop\n  Recommendation Engine for Farmers", 8),
        ("Theme-", " Agriculture, Food Technology & Rural\n  Development", 8),
        ("PS Category-", " Software", 8),
        ("Team Details :", "", 2)
    ]

    for i, (label, val, gap_after) in enumerate(items):
        p = tf_left.paragraphs[0] if i == 0 else tf_left.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(gap_after)
        p.line_spacing = Pt(20)
        
        # Bullet symbol
        run_dot = p.add_run()
        run_dot.text = "• "
        run_dot.font.name = "Arial"
        run_dot.font.size = Pt(17.5)
        run_dot.font.bold = True
        run_dot.font.color.rgb = RGBColor(0, 0, 0)

        # Label
        run_lbl = p.add_run()
        run_lbl.text = label
        run_lbl.font.name = "Arial"
        run_lbl.font.size = Pt(17.5)
        run_lbl.font.bold = True
        run_lbl.font.color.rgb = RGBColor(0, 0, 0)

        # Value
        if val:
            run_val = p.add_run()
            run_val.text = val
            run_val.font.name = "Arial"
            run_val.font.size = Pt(17.5)
            run_val.font.bold = False
            run_val.font.color.rgb = RGBColor(0, 0, 0)

    # Team Members Table / List
    team_members = [
        ("1.", "Sujeet Sandeep Bade", "2025UCP1185"),
        ("2.", "Krishna Maheshwari", "2025UCP1422"),
        ("3.", "Sankirat Kaur", "2025UCP1706"),
        ("4.", "Hemlata Chopda", "2025UCP1343"),
        ("5.", "Arnav Goel", "2025UCP1138"),
        ("6.", "Aman Ali", "2025UCP1488"),
    ]

    for num, name, roll in team_members:
        p = tf_left.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(2)
        p.line_spacing = Pt(18)
        
        run_space = p.add_run()
        run_space.text = f"    {num}  {name:<24}  {roll}"
        run_space.font.name = "Arial"
        run_space.font.size = Pt(15.5)
        run_space.font.color.rgb = RGBColor(0, 0, 0)

    # ----------------------------------------------------
    # SLIDE 2: PROPOSED SOLUTION
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    bg2_path = os.path.join(assets_dir, "user_cleaned_slide_2.png")
    if os.path.exists(bg2_path):
        slide2.shapes.add_picture(bg2_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Slide 2 Title Box
    s2_title_box = slide2.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf2 = s2_title_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "Fasal-Disha (फसल-दिशा)"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = "Times New Roman"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(0, 0, 0)

    p2_sub = tf2.add_paragraph()
    p2_sub.text = "हर खेत को मिले सही दिशा — AI-Powered Net Profit (₹/Acre) Crop Decision Platform"
    p2_sub.alignment = PP_ALIGN.CENTER
    p2_sub.font.name = "Arial"
    p2_sub.font.size = Pt(11)
    p2_sub.font.bold = True
    p2_sub.font.color.rgb = RGBColor(20, 83, 45)

    # Left Column (Core Solution & Unique Differentiators)
    s2_left = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s2_left = s2_left.text_frame
    tf_s2_left.word_wrap = True

    # Box 1: Idea / Proposed Solution
    p = tf_s2_left.paragraphs[0]
    p.text = "❖ IDEA / PROPOSED SOLUTION :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    p = tf_s2_left.add_paragraph()
    p.text = "86.2% of Indian smallholders operate without soil health cards and choose crops by guesswork. Fasal-Disha replaces this with GPS-based soil analysis, government CACP cost data, and real mandi prices — delivering true Net Profit (₹/Acre) forecasts in under 90 seconds."
    p.font.name = "Arial"
    p.font.size = Pt(10.5)
    p.font.color.rgb = RGBColor(30, 41, 59)
    p.space_after = Pt(6)

    s2_bullets = [
        ("Zero-Lab Soil Intelligence: ", "GPS auto-fetches SoilGrids 250m satellite data (pH, clay%, organic carbon) — eliminates the 3-4 week lab test barrier."),
        ("True Net Profit, Not Gross Yield: ", "Uses official CACP A₂+FL cost norms to show real pocket profit per acre — credits tractor (−₹3,500), sprayer (−₹800), pump (−₹600) if owned."),
        ("Crop Rotation Intelligence: ", "Penalizes monoculture (−15% yield penalty) and rewards legume rotation (+12% nitrogen fixation bonus) to prevent soil exhaustion."),
        ("Sowing Window Protection: ", "ICAR-calibrated delay decay (−0.5%/day) warns farmers about invisible yield loss from late planting and auto-suggests contingency crops.")
    ]

    for title, desc in s2_bullets:
        p = tf_s2_left.add_paragraph()
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    # Box 2: Unique Value Propositions
    p = tf_s2_left.add_paragraph()
    p.space_before = Pt(8)
    p.text = "❖ WHAT MAKES IT UNIQUE :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    uvp_bullets = [
        ("Head-to-Head Crop Comparison: ", "Direct ₹/Acre profit comparison between farmer's intended crop and AI optimal recommendation."),
        ("What-If Risk Simulation: ", "Stress-test rainfall deficit (−35%) and mandi price crash (−25%) before spending capital on seeds."),
        ("Multilingual Voice AI: ", "Full voice onboarding and audio playback in Hindi, Marathi, and regional languages for non-literate farmers."),
        ("1-Click Offline Slip: ", "Instant printable A4 advisory slip with complete crop calendar for CSCs, Kisan Kendras & WhatsApp.")
    ]

    for title, desc in uvp_bullets:
        p = tf_s2_left.add_paragraph()
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    # Right Column Box (Decision Ecosystem Diagram Summary)
    s2_right = slide2.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s2_right = s2_right.text_frame
    tf_s2_right.word_wrap = True

    p = tf_s2_right.paragraphs[0]
    p.text = "❖ FASAL-DISHA INTERACTIVE FLOW :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    flow_steps = [
        ("1. 5-Card Farm Profile", "GPS SoilGrids 250m GIS + Texture Cards + Water Source + Owned Machinery + Previous Crop."),
        ("2. Real-Time Mandi & Cost Sync", "14,780+ APMC Mandi price trend points + CACP A₂+FL official operational input costs."),
        ("3. XGBoost Predictive Engine", "District-level yield regression (R²=0.9907) adjusted with ICAR sowing delay curve (−0.5%/day)."),
        ("4. Head-to-Head Profit Matrix", "Farmer's considered choice vs AI recommendation with exact ₹/Acre net gain."),
        ("5. Live What-If Risk Sandbox", "Rainfall sensitivity and mandi crash sliders dynamically recalibrate profit."),
        ("6. 120-Day Crop Calendar & Slip", "Phase-wise farming schedule + 1-Click printable vector advisory slip.")
    ]

    for step, desc in flow_steps:
        p = tf_s2_right.add_paragraph()
        p.space_before = Pt(5)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = step + "\n"
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(71, 85, 105)

    # ----------------------------------------------------
    # SLIDE 3: TECHNICAL APPROACH
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    bg3_path = os.path.join(assets_dir, "user_cleaned_slide_3.png")
    if os.path.exists(bg3_path):
        slide3.shapes.add_picture(bg3_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Slide 3 Title
    s3_title_box = slide3.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf3 = s3_title_box.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "Fasal-Disha Technical Approach"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = "Times New Roman"
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(0, 0, 0)

    p3_sub = tf3.add_paragraph()
    p3_sub.text = "4-Tier Architecture • Zero-Lab GIS Ingestion • XGBoost (R²=0.9907) • 100% CACP Fidelity"
    p3_sub.alignment = PP_ALIGN.CENTER
    p3_sub.font.name = "Arial"
    p3_sub.font.size = Pt(11)
    p3_sub.font.bold = True
    p3_sub.font.color.rgb = RGBColor(20, 83, 45)

    # Slide 3 Left Column (Architecture & Stack)
    s3_left = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s3_left = s3_left.text_frame
    tf_s3_left.word_wrap = True

    p = tf_s3_left.paragraphs[0]
    p.text = "❖ 4-TIER CLOUD-NATIVE ARCHITECTURE :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    tiers = [
        ("Tier 1: Client & Voice UI", "React 18 PWA, TypeScript, Tailwind CSS, Vite, Capacitor 8.5 Native Android Container (< 1.1 MB), Web Speech Recognition, MeitY Indic Speech TTS."),
        ("Tier 2: High-Performance Engine", "FastAPI (Python 3.11 ASGI), Pydantic v2 strict schemas, CACP A₂+FL cost computation engine, ICAR sowing window delay loss calculators."),
        ("Tier 3: Machine Learning Inference", "XGBoostRegressor (R²=0.9907, 95.9% census match) for taluka yields, Agmarknet 5-yr modal pricing forecaster (R²=0.8456)."),
        ("Tier 4: Ground-Truth Repositories", "ISRIC SoilGrids 250m GIS raster layers, MoA Agmarknet APMC mandi archives, CACP Kharif/Rabi price policy norms, PostgreSQL 16 / SQLite fallback.")
    ]

    for title, desc in tiers:
        p = tf_s3_left.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title + ": "
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    # Slide 3 Right Column (Ground Truth Benchmarks)
    s3_right = slide3.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s3_right = s3_right.text_frame
    tf_s3_right.word_wrap = True

    p = tf_s3_right.paragraphs[0]
    p.text = "❖ GROUND-TRUTH ACCURACY BENCHMARKS :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    benchmarks = [
        ("Yield Prediction: ", "XGBoost vs MoA DES Census actuals — 95.9% district yield match (R²=0.9907)."),
        ("Price Forecasting: ", "Model vs Agmarknet APMC actuals — 98.6% wholesale rate accuracy across 14,780+ mandis."),
        ("Cost Calculation: ", "100% CACP A₂+FL fidelity (with verified tractor saving −₹3,500/ac, sprayer −₹800/ac)."),
        ("Sowing Delay Loss: ", "−0.5%/day calibrated from ICAR-CRIDA field trials — 100% ICAR calibration."),
        ("Real Field Proof: ", "Wardha (Maharashtra) farmer, Cotton (monoculture) → Soybean (AI recommendation) = +₹10,370/acre (+72.3% net profit gain).")
    ]

    for title, desc in benchmarks:
        p = tf_s3_right.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.8)
        r2.font.color.rgb = RGBColor(71, 85, 105)

    # ----------------------------------------------------
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    bg4_path = os.path.join(assets_dir, "user_cleaned_slide_4.png")
    if os.path.exists(bg4_path):
        slide4.shapes.add_picture(bg4_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Slide 4 Title
    s4_title_box = slide4.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf4 = s4_title_box.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    p4.text = "Feasibility, Viability & Mitigations"
    p4.alignment = PP_ALIGN.CENTER
    p4.font.name = "Times New Roman"
    p4.font.size = Pt(28)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(0, 0, 0)

    # Slide 4 Left Column
    s4_left = slide4.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s4_left = s4_left.text_frame
    tf_s4_left.word_wrap = True

    p = tf_s4_left.paragraphs[0]
    p.text = "⚖️ FEASIBILITY & VIABILITY ANALYSIS :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    feas_points = [
        ("1. Technical Readiness: ", "Production stack deployed and tested — FastAPI ASGI + React 18 PWA + Capacitor Android (< 1.1 MB bundle). Tested on entry-level 4G devices."),
        ("2. Zero-Lab Removes #1 Barrier: ", "86.2% of Indian smallholders lack lab test kits. Automated 250m GIS SoilGrids satellite ingestion eliminates the 3-4 week wait."),
        ("3. Ultra-Low Overhead: ", "< ₹0.02 operational cost per recommendation query. 100% open-source FOSS licensing."),
        ("4. Deployment Ready: ", "Interoperable open REST APIs ready for integration with PM-Kisan, KVK extension services, and AgriStack national DPI.")
    ]

    for title, desc in feas_points:
        p = tf_s4_left.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    p = tf_s4_left.add_paragraph()
    p.space_before = Pt(8)
    p.text = "🛡️ REAL-WORLD CHALLENGES & MITIGATIONS :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(185, 28, 28)

    challenges = [
        ("Extreme Weather Uncertainty: ", "What-If Risk Sandbox stress-tests rainfall deficit (−35%) and mandi price crashes (−25%) before seed sowing."),
        ("Rural 2G/3G Connectivity: ", "Offline-first PWA caching + 1-Click printable vector A4 advisory slip for CSCs and Kisan Kendras."),
        ("Low Literacy & Language: ", "Full voice navigation and audio narration in regional languages — zero reading barrier.")
    ]

    for title, desc in challenges:
        p = tf_s4_left.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(185, 28, 28)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    # Slide 4 Right Column (Validation Table)
    s4_right = slide4.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s4_right = s4_right.text_frame
    tf_s4_right.word_wrap = True

    p = tf_s4_right.paragraphs[0]
    p.text = "📊 GOVERNMENT DATA GROUND-TRUTH ACCURACY :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    table_data = [
        ("District Yield (qtl/ac)", "MoA DES Census", "95.9% Match"),
        ("Cultivation Cost (₹/ac)", "CACP A₂+FL Norms", "100% Fidelity"),
        ("Mandi Wholesale Rate", "Agmarknet APMC 5-Yr", "98.6% Accuracy"),
        ("Sowing Delay Loss", "ICAR-CRIDA Trials", "100% ICAR Calibrated"),
        ("Soil Spatial Texture", "ISRIC SoilGrids 250m", "100% Geofenced")
    ]

    for param, src, match in table_data:
        p = tf_s4_right.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"• {param:<24} | {src:<18} → "
        r1.font.bold = True
        r1.font.size = Pt(9.8)
        r1.font.color.rgb = RGBColor(30, 41, 59)
        r2 = p.add_run()
        r2.text = match
        r2.font.bold = True
        r2.font.size = Pt(9.8)
        r2.font.color.rgb = RGBColor(21, 128, 61)

    p = tf_s4_right.add_paragraph()
    p.space_before = Pt(10)
    p.text = "★ REAL-WORLD FIELD PROOF ★\nWardha (Maharashtra) • 3.5-Acre Clay Loam • Borewell • Owned Tractor\nCotton (repeat monoculture) → Soybean (AI recommendation)\n₹14,352/ac → ₹24,722/ac = +₹10,370/ac (+72.3% Net Profit Gain)"
    p.font.name = "Arial"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    # ----------------------------------------------------
    # SLIDE 5: IMPACT AND BENEFITS
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    bg5_path = os.path.join(assets_dir, "user_cleaned_slide_5.png")
    if os.path.exists(bg5_path):
        slide5.shapes.add_picture(bg5_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Slide 5 Title
    s5_title_box = slide5.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf5 = s5_title_box.text_frame
    tf5.word_wrap = True
    p5 = tf5.paragraphs[0]
    p5.text = "Transformative Impact & Measurable Benefits"
    p5.alignment = PP_ALIGN.CENTER
    p5.font.name = "Times New Roman"
    p5.font.size = Pt(28)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(0, 0, 0)

    # Slide 5 Left Column (Impact)
    s5_left = slide5.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s5_left = s5_left.text_frame
    tf_s5_left.word_wrap = True

    p = tf_s5_left.paragraphs[0]
    p.text = "👥 TARGET AUDIENCE IMPACT (14.6 Cr Indian Farmers) :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    impact_points = [
        ("Vulnerable Demographic Focus: ", "Equips 120M+ smallholders (< 2 Ha) with institutional-grade ₹/Acre crop economics — zero lab tests, zero fees."),
        ("Elimination of Guesswork: ", "Replaces backward memory-based planting with predictive Net In-Hand Profit (₹/Acre) accounting."),
        ("Zero-Surprise Farming: ", "Not just 'grow wheat' — full economics + risk simulator + 120-day milestone calendar + printable advisory slip.")
    ]

    for title, desc in impact_points:
        p = tf_s5_left.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.2)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.8)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    p = tf_s5_left.add_paragraph()
    p.space_before = Pt(8)
    p.text = "🌐 4-DIMENSIONAL NATIONAL BENEFITS :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    benefits = [
        ("💰 Economic: ", "+₹5,000 to ₹12,000/ac net profit increase. Owned machinery saves up to ₹4,900/ac."),
        ("🌿 Ecological: ", "+12% biological nitrogen fixation from legume rotation, reversing soil monoculture exhaustion."),
        ("🗣️ Social Inclusion: ", "Regional voice narration bridges literacy barrier; 1-Click PDF slip enables offline sharing."),
        ("🇮🇳 Policy Alignment: ", "100% FOSS architecture ready for PM-Kisan, KVK extension & AgriStack DPI.")
    ]

    for title, desc in benefits:
        p = tf_s5_left.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10.2)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.8)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    # Slide 5 Right Column (Before vs After)
    s5_right = slide5.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s5_right = s5_right.text_frame
    tf_s5_right.word_wrap = True

    p = tf_s5_right.paragraphs[0]
    p.text = "⚖️ BEFORE VS AFTER PARADIGM SHIFT :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    shifts = [
        ("Traditional Guesswork", "Fasal-Disha AI Precision", "Chasing gross harvest yield → True Net In-Hand Profit (₹/Acre)"),
        ("3-4 Week Soil Lab Wait", "Instant GIS Satellite Ingestion", "86.2% unserved farmers → 0-second GPS SoilGrids profiling"),
        ("Repetitive Monoculture", "Smart Crop Rotation Engine", "Soil nitrogen exhaustion → +12% biological nitrogen replenishment"),
        ("Blind Climate & Market Risk", "Live Risk Sandbox Simulator", "Debt trap from price drops → Pre-sowing shock testing (−35% rain / −25% price)"),
        ("Text-Heavy Complex Dashboards", "Voice AI + 1-Click Printable Slip", "Literacy barrier → 100% spoken regional voice + WhatsApp PDF slip")
    ]

    for b, a, note in shifts:
        p = tf_s5_right.add_paragraph()
        p.space_before = Pt(5)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = f"❌ {b}\n"
        r1.font.bold = True
        r1.font.size = Pt(9.8)
        r1.font.color.rgb = RGBColor(185, 28, 28)
        r2 = p.add_run()
        r2.text = f"✅ {a} ({note})\n"
        r2.font.bold = True
        r2.font.size = Pt(9.8)
        r2.font.color.rgb = RGBColor(21, 128, 61)

    # ----------------------------------------------------
    # SLIDE 6: RESEARCH AND REFERENCES
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    bg6_path = os.path.join(assets_dir, "user_cleaned_slide_6.png")
    if os.path.exists(bg6_path):
        slide6.shapes.add_picture(bg6_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Slide 6 Title
    s6_title_box = slide6.shapes.add_textbox(Inches(3.0), Inches(0.2), Inches(7.333), Inches(0.9))
    tf6 = s6_title_box.text_frame
    tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "Research, References & Prototype Deliverables"
    p6.alignment = PP_ALIGN.CENTER
    p6.font.name = "Times New Roman"
    p6.font.size = Pt(28)
    p6.font.bold = True
    p6.font.color.rgb = RGBColor(0, 0, 0)

    # Slide 6 Left Column
    s6_left = slide6.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.6))
    tf_s6_left = s6_left.text_frame
    tf_s6_left.word_wrap = True

    p = tf_s6_left.paragraphs[0]
    p.text = "🏛️ OFFICIAL GOVERNMENT & SCIENTIFIC CITATIONS :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 83, 45)

    citations = [
        ("CACP (Commission for Agricultural Costs & Prices): ", "Price Policy for Kharif/Rabi Crops (2020-2025) — State-wise itemized A₂+FL operational cultivation cost benchmarks [cacp.dacnet.nic.in]."),
        ("Agmarknet (Ministry of Agriculture): ", "14,780+ daily APMC mandi wholesale records across 5 years for seasonal modal price forecasting [agmarknet.gov.in]."),
        ("ICAR-CRIDA: ", "Agro-climatic crop calendars and daily yield decay curves (≈0.5%/day) for late sowing contingency planning [icar.org.in]."),
        ("ISRIC World Soil Information (SoilGrids 250m): ", "Global 250m gridded rasters for Organic Carbon, pH, Bulk Density, Sand/Clay fractions [soilgrids.org]."),
        ("Chen & Guestrin (2016): ", "XGBoost: A Scalable Tree Boosting System (ACM SIGKDD).")
    ]

    for title, desc in citations:
        p = tf_s6_left.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(10)
        r1.font.color.rgb = RGBColor(20, 83, 45)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.5)
        r2.font.color.rgb = RGBColor(51, 65, 85)

    # Slide 6 Right Column (Verified Working Prototype)
    s6_right = slide6.shapes.add_textbox(Inches(7.6), Inches(1.3), Inches(5.2), Inches(5.6))
    tf_s6_right = s6_right.text_frame
    tf_s6_right.word_wrap = True

    p = tf_s6_right.paragraphs[0]
    p.text = "🚀 VERIFIED WORKING PROTOTYPE & DELIVERABLES :"
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138)

    deliverables = [
        ("💻 GitHub Repository (Public): ", "github.com/sujeetaionly/agri-decide-sih (100% Tested Production Stack, 12/12 passing pytest suite)."),
        ("🎬 2-Minute Demo Video: ", "Full Android Mobile App Walkthrough demonstrating Zero-Lab Onboarding, What-If Risk Sandbox, and 1-Click PDF Advisory Slip."),
        ("📱 Native Android App: ", "Capacitor 8.5 Native APK container (< 1.1 MB) + Offline-first PWA."),
        ("🗣️ Local Language Voice: ", "MeitY Digital India BHASHINI Speech TTS engine + browser Web Speech recognition.")
    ]

    for title, desc in deliverables:
        p = tf_s6_right.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(3)
        r1 = p.add_run()
        r1.text = "• " + title + "\n"
        r1.font.bold = True
        r1.font.size = Pt(10.5)
        r1.font.color.rgb = RGBColor(30, 58, 138)
        r2 = p.add_run()
        r2.text = desc
        r2.font.bold = False
        r2.font.size = Pt(9.8)
        r2.font.color.rgb = RGBColor(71, 85, 105)

    # ----------------------------------------------------
    # SLIDE 7: INTERACTIVE USER EXPERIENCE & WORKFLOW
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    bg7_path = os.path.join(assets_dir, "user_cleaned_slide_2.png")
    if os.path.exists(bg7_path):
        slide7.shapes.add_picture(bg7_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    # Team Oval Text (Top Left)
    oval_box = slide7.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(1.8), Inches(1.0))
    tf_oval = oval_box.text_frame
    tf_oval.word_wrap = True
    p_oval = tf_oval.paragraphs[0]
    p_oval.text = "Almost AC"
    p_oval.alignment = PP_ALIGN.CENTER
    p_oval.font.name = "Arial"
    p_oval.font.size = Pt(20)
    p_oval.font.bold = True
    p_oval.font.color.rgb = RGBColor(0, 0, 0)

    # Slide 7 Title - Centered strictly between Team Oval and SIH logo
    s7_title_box = slide7.shapes.add_textbox(Inches(2.4), Inches(0.18), Inches(7.8), Inches(1.1))
    tf7 = s7_title_box.text_frame
    tf7.word_wrap = True
    p7 = tf7.paragraphs[0]
    p7.text = "Screenshots of Real Working Prototype"
    p7.alignment = PP_ALIGN.CENTER
    p7.font.name = "Times New Roman"
    p7.font.size = Pt(25)
    p7.font.bold = True
    p7.font.color.rgb = RGBColor(0, 0, 0)

    p7_sub = tf7.add_paragraph()
    p7_sub.text = "12+ Indian Regional Languages ➔ 250m GIS Soil Profiling ➔ Head-to-Head Profit Engine ➔ What-If Risk Sandbox ➔ 120-Day Action Plan"
    p7_sub.alignment = PP_ALIGN.CENTER
    p7_sub.font.name = "Arial"
    p7_sub.font.size = Pt(9.5)
    p7_sub.font.bold = True
    p7_sub.font.color.rgb = RGBColor(20, 83, 45)

    # 5 Flow Step Columns - Starting strictly below Y = 1.35 Inches
    flow_steps_data = [
        (
            "flow_step_1_home.png",
            "STEP 1: Home & Voice AI",
            "All Regional Languages [BHASHINI]:\nVoice guidance in 12+ Indian languages, live Mandi ticker & 1-tap 2-min assessment."
        ),
        (
            "flow_step_2_soil.png",
            "STEP 2: Zero-Lab GIS Profile",
            "Zero-Lab Satellite GIS [250m Grid]:\nAuto-fetches SoilGrids 250m texture/pH. Credits owned machinery & rotation."
        ),
        (
            "flow_step_3_comparison.png",
            "STEP 3: Head-to-Head Decision",
            "Intended vs AI Best [+₹15,016/Ac]:\nFarmer Cotton vs AI Moong. Itemized CACP A₂+FL cost & net profit delta."
        ),
        (
            "flow_step_4_whatif.png",
            "STEP 4: What-If Risk Sandbox",
            "Pre-Sowing Stress-Test [Real-Time]:\nSliders for −30% rain & −25% mandi crash. Dynamic recalibration (₹19,060/ac)."
        ),
        (
            "flow_step_5_calendar.png",
            "STEP 5: 120-Day Action Plan",
            "Operational Guidance [1-Click PDF]:\nDay-by-day seed, weeding & harvest plan + 1-Click offline printable QR slip."
        ),
    ]

    col_width = Inches(2.35)
    col_gap = Inches(0.2)
    start_left = Inches(0.38)

    for idx, (img_name, step_header, step_desc) in enumerate(flow_steps_data):
        curr_left = start_left + idx * (col_width + col_gap)

        # Step Header Box (Below Oval / Logo Clearance Zone)
        hdr_box = slide7.shapes.add_textbox(curr_left, Inches(1.36), col_width, Inches(0.36))
        tf_hdr = hdr_box.text_frame
        tf_hdr.word_wrap = True
        p_hdr = tf_hdr.paragraphs[0]
        p_hdr.text = step_header
        p_hdr.alignment = PP_ALIGN.CENTER
        p_hdr.font.name = "Arial"
        p_hdr.font.size = Pt(10.5)
        p_hdr.font.bold = True
        p_hdr.font.color.rgb = RGBColor(21, 128, 61)

        # Phone Mockup Image (Border-free)
        img_path = os.path.join(assets_dir, img_name)
        if os.path.exists(img_path):
            slide7.shapes.add_picture(img_path, curr_left, Inches(1.75), width=col_width, height=Inches(4.55))

        # Step Caption Box
        desc_box = slide7.shapes.add_textbox(curr_left, Inches(6.35), col_width, Inches(0.85))
        tf_desc = desc_box.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = step_desc
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.font.name = "Arial"
        p_desc.font.size = Pt(9)
        p_desc.font.color.rgb = RGBColor(51, 65, 85)

    try:
        prs.save(output_path)
        print(f"Presentation saved successfully to: {output_path}")
    except PermissionError:
        alt_path = os.path.abspath("presentation/Fasal_Disha_SIH_2025_7slides.pptx")
        prs.save(alt_path)
        print(f"File locked by PowerPoint. Saved to alternative path: {alt_path}")

if __name__ == "__main__":
    create_deck_pptx()

